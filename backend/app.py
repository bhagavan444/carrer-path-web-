from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import os
import re
import uuid
import json
import requests
import logging
import random
from concurrent.futures import ThreadPoolExecutor
import pdfplumber
import docx2txt
from collections import Counter
from datetime import datetime
from docx import Document
from pptx import Presentation
from fpdf import FPDF
import pandas as pd
import threading

# -------------------- Setup --------------------
app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}})

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Thread-safe chat session storage
chat_sessions = {}
chat_sessions_lock = threading.Lock()

# Thread pool for concurrent API calls
executor = ThreadPoolExecutor(max_workers=5)

# -------------------- Config --------------------
career_options = {
    "python": ["Data Scientist", "Python Developer", "AI/ML Engineer"],
    "javascript": ["Full Stack Developer", "React Developer", "Frontend Engineer"],
    "design": ["UI/UX Designer", "Product Designer", "Graphic Designer"],
    "ai": ["AI Engineer", "Research Scientist", "NLP Engineer"],
    "cybersecurity": ["Cybersecurity Analyst", "Ethical Hacker", "SOC Analyst"],
    "java": ["Backend Developer", "Spring Boot Engineer", "Android Developer"],
}

certification_suggestions = {
    "Data Scientist": ["IBM Data Science", "Google Advanced ML", "Coursera AI for Everyone"],
    "UI/UX Designer": ["Google UX Design", "Adobe Certified Designer"],
    "Cybersecurity Analyst": ["CompTIA Security+", "Certified Ethical Hacker"],
    "Python Developer": ["Python for Everybody", "Google IT Automation with Python"],
    "AI Engineer": ["DeepLearning.ai Specialization", "OpenAI Bootcamp"],
}

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "AIzaSyA3uC2IDU_Rb8VoJ-k2lGILQROc7j0SgNU")
GEMINI_API_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent"
UPLOAD_API_URL = "https://generativelanguage.googleapis.com/upload/v1beta/files"

os.makedirs("temp", exist_ok=True)
os.makedirs("downloads", exist_ok=True)

# -------------------- Helpers --------------------
def extract_resume_text(filepath):
    """Extract plain text from PDF/DOCX/TXT resumes"""
    try:
        if filepath.endswith(".pdf"):
            with pdfplumber.open(filepath) as pdf:
                return "\n".join(page.extract_text() or "" for page in pdf.pages)
        elif filepath.endswith(".docx"):
            return docx2txt.process(filepath)
        else:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception as e:
        logger.error(f"Error extracting resume text: {e}")
        raise

def extract_keywords(text):
    words = re.findall(r"\b[a-zA-Z]+\b", text.lower())
    freq = Counter(words)
    matched = [key for key in career_options if key in freq]
    return matched, freq

def extract_resume_summary(text):
    text = text.lower()
    def find_section(section):
        match = re.search(rf"{section}[\s\S]{{0,300}}", text)
        return match.group(0).strip() if match else "Not found"
    return {
        "education": find_section("education"),
        "experience": find_section("experience"),
        "projects": find_section("project"),
    }

def suggest_skill_gaps(found, roles):
    expected = {
        "Data Scientist": ["python", "numpy", "pandas", "ml", "statistics"],
        "UI/UX Designer": ["figma", "wireframes", "user testing", "adobe"],
        "Cybersecurity Analyst": ["network", "firewall", "threat", "incident"],
        "Python Developer": ["flask", "django", "oop", "algorithms"],
        "AI Engineer": ["deep learning", "tensorflow", "nlp", "neural"],
    }
    gaps = []
    for role in roles:
        for word in expected.get(role, []):
            if word not in found:
                gaps.append(word)
    return list(set(gaps))

# -------------------- Gemini API Call --------------------
def call_gemini_api(prompt):
    try:
        headers = {"Content-Type": "application/json", "x-goog-api-key": GEMINI_API_KEY}
        body = {"contents": [{"role": "user", "parts": [{"text": prompt}]}]}
        response = requests.post(GEMINI_API_URL, headers=headers, json=body, timeout=10)
        response.raise_for_status()
        return response.json()
    except Exception as e:
        logger.error(f"Gemini API error: {e}")
        raise

# -------------------- Career Path Predictor --------------------
@app.route("/predict", methods=["POST"])
def predict():
    preferred = request.form.get("preferredDomain", "").lower()
    use_ai = request.form.get("useAI") == "true"
    uploaded_file = request.files.get("resume")

    if not uploaded_file:
        return jsonify({"error": "No resume uploaded."}), 400

    filepath = os.path.join("temp", uploaded_file.filename)
    try:
        uploaded_file.save(filepath)
        resume_text = extract_resume_text(filepath)
    except Exception as e:
        return jsonify({"error": f"Resume processing error: {str(e)}"}), 500
    finally:
        if os.path.exists(filepath):
            os.remove(filepath)

    matched_keywords, _ = extract_keywords(resume_text)
    summary = extract_resume_summary(resume_text)

    roles = []
    if use_ai or not preferred:
        for kw in matched_keywords:
            roles.extend(career_options.get(kw, []))
    else:
        for key, vals in career_options.items():
            if preferred in key or preferred in vals[0].lower():
                roles = vals
                break

    if not roles:
        roles = ["General Software Engineer", "Tech Associate", "Junior Developer"]

    certs = []
    for r in roles:
        certs.extend(certification_suggestions.get(r, []))

    skill_gaps = suggest_skill_gaps(matched_keywords, roles)
    score = min(95, max(65, 65 + len(matched_keywords) * 2))

    return jsonify({
        "roles": list(set(roles)),
        "skills": matched_keywords or ["teamwork", "problem-solving"],
        "score": score,
        "certifications": list(set(certs)),
        "skillGaps": skill_gaps,
        "summary": summary,
        "suggestions": "Use action verbs, quantify impact, and highlight relevant projects."
    })

# -------------------- Chat Endpoints --------------------
@app.route("/api/chats", methods=["GET"])
def get_chats():
    with chat_sessions_lock:
        return jsonify({"sessions": [
            {"_id": cid, "messages": s["messages"], "title": s.get("title", "Untitled Chat")}
            for cid, s in chat_sessions.items()
        ]})
@app.route("/api/chats/<chat_id>", methods=["PATCH"])
def update_chat(chat_id):
    data = request.get_json()
    new_title = data.get("title")
    if not new_title:
        return jsonify({"error": "Title is required"}), 400
    with chat_lock:
        if chat_id in chat_sessions:
            chat_sessions[chat_id]["title"] = new_title
            return jsonify({"message": "Title updated successfully", "title": new_title})
    return jsonify({"error": "Chat session not found"}), 404

@app.route("/api/chats/<chat_id>", methods=["GET"])
def get_chat(chat_id):
    with chat_sessions_lock:
        return jsonify(chat_sessions.get(chat_id) or {"error": "Not found"}), (200 if chat_id in chat_sessions else 404)

@app.route("/api/chats", methods=["DELETE"])
def clear_chats():
    with chat_sessions_lock:
        chat_sessions.clear()
    return jsonify({"message": "All chats cleared"})

@app.route("/api/chat", methods=["POST"])
def chat():
    try:
        if request.content_type.startswith("multipart/form-data"):
            user_input = request.form.get("message", "")
            chat_id = request.form.get("chat_id", str(uuid.uuid4()))
            file = request.files.get("file")
        else:
            data = request.get_json() or {}
            user_input = data.get("message", "")
            chat_id = data.get("chat_id", str(uuid.uuid4()))
            file = None

        if not user_input and not file:
            return jsonify({"reply": "⚠️ No input provided"}), 400

        filename, file_uri, mime_type = None, None, None
        if file:
            filename = file.filename.lower()
            if filename.endswith(".docx"):
                text = "\n".join([p.text for p in Document(file).paragraphs])
                user_input += f"\n\nFile Content: {text}"
            elif filename.endswith(".pptx"):
                text = ""
                prs = Presentation(file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                user_input += f"\n\nFile Content: {text}"
            elif filename.endswith((".xlsx", ".xls")):
                df = pd.read_excel(file)
                user_input += f"\n\nFile Content: {df.to_string()}"
            else:
                return jsonify({"reply": "⚠️ Unsupported file type"}), 400

        # Prompt Gemini
        system = """You are a helpful assistant.
If input looks like a resume, return JSON {"type":"ats","score":0-100,"feedback":"..."}.
If user asks PPT, return {"type":"ppt","slides":[...]}.
If PDF editing, return {"type":"pdf","content":"..."}.
Otherwise reply as plain text."""
        prompt = system + "\n\n" + user_input

        # Execute Gemini API call in thread pool
        future = executor.submit(call_gemini_api, prompt)
        result = future.result()

        reply = result.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "")
        download_url = None

        try:
            parsed = json.loads(reply)
            if parsed.get("type") == "ppt":
                prs = Presentation()
                for slide in parsed["slides"]:
                    layout = prs.slide_layouts[1]
                    s = prs.slides.add_slide(layout)
                    s.shapes.title.text = slide.get("title", "Slide")
                    tf = s.placeholders[1].text_frame
                    for l in slide.get("content", "").split("\n"):
                        tf.add_paragraph().text = l
                fname = f"ppt_{uuid.uuid4()}.pptx"
                prs.save(f"downloads/{fname}")
                reply, download_url = "PPT generated ✅", f"/download/{fname}"
            elif parsed.get("type") == "pdf":
                pdf = FPDF()
                pdf.add_page()
                pdf.set_font("Arial", size=12)
                pdf.multi_cell(0, 10, parsed.get("content", ""))
                fname = f"pdf_{uuid.uuid4()}.pdf"
                pdf.output(f"downloads/{fname}")
                reply, download_url = "PDF generated ✅", f"/download/{fname}"
            elif parsed.get("type") == "ats":
                reply = f"ATS Score: {parsed.get('score',0)}/100\n{parsed.get('feedback','')}"
        except Exception:
            pass

        # Store chat history
        with chat_sessions_lock:
            if chat_id not in chat_sessions:
                chat_sessions[chat_id] = {"messages": [], "title": user_input[:20] + "..."}
            ts = datetime.now().strftime("%I:%M %p")
            chat_sessions[chat_id]["messages"].append({"id": str(uuid.uuid4()), "role": "user", "message": user_input, "time": ts})
            chat_sessions[chat_id]["messages"].append({"id": str(uuid.uuid4()), "role": "assistant", "reply": reply, "time": ts})

        resp = {"reply": reply, "chat_id": chat_id}
        if download_url: resp["download_url"] = download_url
        return jsonify(resp)

    except Exception as e:
        logger.error(f"Chat error: {e}")
        return jsonify({"reply": "⚠️ Error processing message"}), 500

# -------------------- File Download --------------------
@app.route("/download/<fname>")
def download(fname):
    try:
        return send_from_directory("downloads", fname, as_attachment=True)
    except FileNotFoundError:
        return jsonify({"error": "File not found"}), 404

# -------------------- Run --------------------
if __name__ == "__main__":
    try:
        app.run(debug=True, port=5000, threaded=True)
    finally:
        executor.shutdown(wait=True)